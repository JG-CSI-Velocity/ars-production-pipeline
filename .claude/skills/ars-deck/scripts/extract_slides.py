#!/usr/bin/env python3
"""Extract per-slide text + embedded chart images from an ARS deck.

LibreOffice (full-slide render) usually isn't installed on the dev Mac, so to
actually *see* what a slide looks like, pull its embedded chart PNGs and read
them. Also prints an inventory (title + picture count) so you can map slide
numbers to sections.

    python extract_slides.py deck.pptx                 # inventory + all images
    python extract_slides.py deck.pptx --slides 33-43   # just a range
    python extract_slides.py deck.pptx --slides 12,53,167 --out /tmp/review
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE


def parse_slides(spec: str, total: int) -> list[int]:
    if not spec:
        return list(range(1, total + 1))
    out: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            lo, hi = part.split("-")
            out.extend(range(int(lo), int(hi) + 1))
        elif part:
            out.append(int(part))
    return [s for s in out if 1 <= s <= total]


def title_of(slide) -> str:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0][:70]
    return "(no text)"


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx")
    ap.add_argument("--slides", default="", help="e.g. 33-43 or 12,53,167")
    ap.add_argument("--out", default="/tmp/ars_slides")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    out = Path(args.out)
    out.mkdir(parents=True, exist_ok=True)
    targets = parse_slides(args.slides, len(prs.slides))

    print(f"{args.pptx}: {len(prs.slides)} slides | extracting {len(targets)} -> {out}\n")
    for idx in targets:
        slide = prs.slides[idx - 1]
        n_img = 0
        for sh in slide.shapes:
            if sh.shape_type == PICTURE:
                n_img += 1
                (out / f"s{idx:03d}_{n_img}.png").write_bytes(sh.image.blob)
        texts = [
            sh.text_frame.text.replace("\n", " ").strip()[:80]
            for sh in slide.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()
        ]
        print(f"  s{idx:>3} | {n_img} img | {title_of(slide)}")
        for t in texts[1:4]:
            print(f"        - {t}")
    print(f"\nRead the PNGs in {out} to judge the visuals.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
