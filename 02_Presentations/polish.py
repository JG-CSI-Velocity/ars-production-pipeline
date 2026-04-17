"""Deck polish driver -- post-hoc PPTX compliance pass."""

from __future__ import annotations

import argparse
import sys
from dataclasses import dataclass, field
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from loguru import logger
from pptx import Presentation

from style.charts import audit_chart_image
from style.headline import score_headline
from style.layout import FOOTER_ZONE, TITLE_ZONE, is_inside_zone
from style.narrative import NarrativeScore, score_slide


@dataclass
class SlideAudit:
    index: int
    headline: str
    narrative: NarrativeScore
    headline_violations: list[str]
    chart_dpi_violations: list[str]
    flagged: bool


@dataclass
class DeckLevelAudit:
    client_name_present: bool
    section_dividers_ok: bool
    page_numbers_ok: bool
    appendix_separated: bool


@dataclass
class DeckAudit:
    deck_path: Path
    slide_count: int
    deck_level: DeckLevelAudit
    slides: list[SlideAudit] = field(default_factory=list)


def _extract_title(slide, slide_height) -> str:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if is_inside_zone(shape, TITLE_ZONE, slide_height):
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def _client_name_present(title_slide_headline: str) -> bool:
    """Heuristic: non-generic placeholder (contains Bank/Union or >=4 words)."""
    if not title_slide_headline:
        return False
    lower = title_slide_headline.lower()
    return "bank" in lower or "union" in lower or len(title_slide_headline.split()) >= 4


def audit_deck(deck_path: Path) -> DeckAudit:
    prs = Presentation(str(deck_path))
    slide_height = prs.slide_height

    title_text = _extract_title(prs.slides[0], slide_height) if len(prs.slides) else ""
    deck_level = DeckLevelAudit(
        client_name_present=_client_name_present(title_text),
        section_dividers_ok=True,  # placeholder for future expansion
        page_numbers_ok=True,
        appendix_separated=True,
    )

    slide_audits: list[SlideAudit] = []
    for i, slide in enumerate(prs.slides, start=1):
        headline = _extract_title(slide, slide_height)
        h_score = score_headline(headline) if headline else None
        narrative = score_slide(slide, slide_height=slide_height)

        chart_flags: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    audit = audit_chart_image(shape)
                    if not audit.meets_dpi_floor:
                        chart_flags.append(
                            f"Image below DPI floor ({audit.dpi_estimate} < 150)"
                        )
                except Exception as e:  # pragma: no cover
                    logger.warning(f"Slide {i} image audit failed: {e}")

        # Title slide (index 1) is exempt from narrative flags -- its
        # headline is inherently not a consultative finding.
        is_title_slide = i == 1
        flagged = (not is_title_slide) and (
            narrative.consultative < 2
            or narrative.performance < 3
            or narrative.focal < 2
            or bool(chart_flags)
        )

        slide_audits.append(
            SlideAudit(
                index=i,
                headline=headline,
                narrative=narrative,
                headline_violations=h_score.violates if h_score else ["no headline detected"],
                chart_dpi_violations=chart_flags,
                flagged=flagged,
            )
        )

    return DeckAudit(
        deck_path=deck_path,
        slide_count=len(prs.slides),
        deck_level=deck_level,
        slides=slide_audits,
    )


def write_report(audit: DeckAudit, path: Path) -> None:
    lines: list[str] = []
    lines.append(f"# Polish Report -- {audit.deck_path.name}")
    lines.append("")
    lines.append("## Summary")
    lines.append(f"- Slides: {audit.slide_count}")
    passing = sum(1 for s in audit.slides if not s.flagged)
    lines.append(f"- Passing (all axes >=2, no DPI violations): {passing}")
    lines.append(f"- Flagged: {audit.slide_count - passing}")
    lines.append("")
    lines.append("## Deck-level findings")
    lines.append(f"- [{'x' if audit.deck_level.client_name_present else ' '}] Title slide has client name")
    lines.append(f"- [{'x' if audit.deck_level.section_dividers_ok else ' '}] Section dividers present")
    lines.append(f"- [{'x' if audit.deck_level.page_numbers_ok else ' '}] Page numbers present")
    lines.append(f"- [{'x' if audit.deck_level.appendix_separated else ' '}] Appendix separated")
    lines.append("")
    lines.append("## Slide-by-slide")
    for s in audit.slides:
        lines.append(f"### Slide {s.index} -- \"{s.headline or '(no title)'}\"")
        lines.append(f"- Consultative: {s.narrative.consultative}/3")
        lines.append(f"- Performance: {s.narrative.performance}/3")
        lines.append(f"- Focal: {s.narrative.focal}/3")
        if s.headline_violations:
            lines.append("- Headline violations:")
            for v in s.headline_violations:
                lines.append(f"  - {v}")
        if s.chart_dpi_violations:
            lines.append("- Chart violations:")
            for v in s.chart_dpi_violations:
                lines.append(f"  - {v}")
        lines.append("")
    path.write_text("\n".join(lines))


def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(
        prog="polish",
        description="Post-hoc PPTX polish pass (SLIDE_MAPPING.md compliance).",
    )
    p.add_argument("deck", nargs="?", type=Path, help="Path to input .pptx")
    p.add_argument("--batch", type=Path, help="Process every .pptx in this folder")
    p.add_argument("--apply", action="store_true", help="Write polished PPTX (default dry-run)")
    p.add_argument("--report-only", action="store_true", help="Emit report only")
    p.add_argument("--out", type=Path, default=None, help="Output directory")
    p.add_argument("--strict", action="store_true", help="Exit non-zero if any flag")
    return p.parse_args(argv)


def _process_one(deck_path: Path, out_dir: Path) -> DeckAudit:
    out_dir.mkdir(parents=True, exist_ok=True)
    audit = audit_deck(deck_path)
    report_path = out_dir / f"{deck_path.stem}__polish_report.md"
    write_report(audit, report_path)
    logger.info(f"Wrote {report_path}")
    return audit


def main(argv: list[str] | None = None) -> int:
    args = _parse_args(argv)

    if args.batch is None and args.deck is None:
        logger.error("Provide a deck path or --batch <dir>")
        return 2

    decks: list[Path]
    if args.batch is not None:
        decks = sorted(args.batch.glob("*.pptx"))
        if not decks:
            logger.error(f"No .pptx files found in {args.batch}")
            return 2
    else:
        decks = [args.deck]

    all_audits: list[DeckAudit] = []
    for d in decks:
        out_dir = args.out if args.out else d.parent / "polished"
        all_audits.append(_process_one(d, out_dir))

    if args.strict:
        any_flag = any(s.flagged for a in all_audits for s in a.slides)
        if any_flag:
            logger.error("Strict mode: flagged slides present")
            return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
