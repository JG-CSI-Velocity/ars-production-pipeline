"""
Phase 18.1 / 18.2 / 18.3 — deck_builder.py reference patch
==========================================================

Target file: 01_Analysis/00-Scripts/output/deck_builder.py

This is a *reference* patch — copy the methods below into the existing
DeckBuilder class. Each block is annotated with WHERE, WHAT, and WHY.

Applied against deck_builder.py as of branch `pipeline-improvement`:
  - SlideContent dataclass (line ~73): has `kpis`, `notes_text` fields. Good.
  - DeckBuilder._add_slide (line ~196): single-arg today, needs to accept
    slide_number/client_name/month.
  - DeckBuilder.build (line ~158): single-arg signature today.
  - DeckBuilder._build_section_slide (line ~428): currently uses default
    placeholder styling on LAYOUT_TITLE.
  - build_deck() at module level (line ~1730): calls
    `builder.build(all_slides, str(pptx_path))` — needs client_name + month.
  - _section_divider() helper inside build_deck (line ~1824): currently emits
    `f"{title}\\n{subtitle}"` with `_DEFAULT_DIVIDER_LAYOUT` (= LAYOUT_TITLE).
    Needs the 3-line "number\\nlabel\\nlead-in" format.

WORK PC NOTE: do not apply this patch by hand. The dev applies it, commits,
operator pulls via GitHub Desktop. See docs/deck/phase-18-operator-guide.md.
"""

# =============================================================================
# IMPORTS (add to the top of deck_builder.py if not already present)
# =============================================================================
# from pptx.enum.shapes import MSO_SHAPE
# (pptx.util.Inches, Pt, RGBColor, PP_ALIGN already imported)


# =============================================================================
# Section number + lead-in constants
# Add near the top of deck_builder.py, after _SECTION_LABELS (line ~1331).
# =============================================================================

_SECTION_NUMBERS = {
    "overview": "01",
    "dctr": "02",
    "rege": "03",
    "attrition": "04",
    "mailer": "05",
    "transaction": "06",
    "ics": "07",
    "value": "08",
    "insights": "09",
}

_SECTION_LEAD_INS = {
    "overview": "Portfolio composition, eligibility, and program scope.",
    "dctr": "Current penetration, where opportunity sits, and what closing the gap is worth.",
    "rege": "Opt-in rates, branch variation, and revenue impact of Reg E enrollment.",
    "attrition": "Who is leaving, what it costs, and where retention efforts should focus.",
    "mailer": "Campaign response rates, cohort lift, and mailer program ROI.",
    "transaction": "Spending patterns, merchant concentration, and engagement signals.",
    "ics": "Invitation Checking System acquisition channels and conversion.",
    "value": "Revenue attribution — what a debit card and Reg E opt-in are worth.",
    "insights": "Synthesis, recommendations, and quantified action plan.",
}


# =============================================================================
# 18.1 — Callout box
# Add as a new method on DeckBuilder, after _add_fitted_picture.
# Called from _add_slide for content slide types.
# SLIDE_DESIGN.md §7.
# =============================================================================

def _add_callout_box(self, slide, kpis, position="bottom_right"):
    """Render a callout box with a hero KPI + sub-label.

    kpis: dict of {label: value}. First non-subtitle entry is the hero number.
    position: "bottom_right" (default) or "bottom_left".
    """
    if not kpis:
        return

    # Pick hero + optional second KPI
    hero_label = hero_value = sub_label = None
    for label, value in kpis.items():
        if label.lower() in ("subtitle", "title"):
            continue
        if hero_value is None:
            hero_label, hero_value = label, str(value)
        elif sub_label is None:
            sub_label = f"{label}: {value}"

    if hero_value is None:
        return

    from pptx.enum.shapes import MSO_SHAPE

    box_w, box_h = Inches(3.8), Inches(1.4)
    slide_w, slide_h = Inches(13.33), Inches(7.5)
    footer_clearance = Inches(0.75)  # space reserved for footer band

    if position == "bottom_left":
        box_left = Inches(0.86)
    else:
        box_left = slide_w - box_w - Inches(0.5)
    box_top = slide_h - box_h - footer_clearance

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_w, box_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xED, 0xFA, 0xF9)
    shape.line.color.rgb = RGBColor(0x0D, 0x94, 0x88)
    shape.line.width = Pt(1.5)
    try:
        shape.adjustments[0] = 0.15
    except Exception:
        pass

    def _textbox(top_offset, height, text, size, bold, color, italic=False):
        tb = slide.shapes.add_textbox(
            box_left + Inches(0.2), box_top + top_offset,
            box_w - Inches(0.4), height,
        )
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT

    _textbox(Inches(0.1), Inches(0.6), hero_value, 32, True, RGBColor(0x0D, 0x94, 0x88))
    _textbox(Inches(0.7), Inches(0.3), hero_label, 14, True, RGBColor(0x1E, 0x3D, 0x59))
    if sub_label:
        _textbox(Inches(1.0), Inches(0.3), sub_label, 12, False, RGBColor(0x33, 0x33, 0x33))


# =============================================================================
# 18.2 — Footer band
# Add as a new method on DeckBuilder, after _add_callout_box.
# Called from _add_slide for all content slides (not title, not section).
# SLIDE_DESIGN.md §8.
# =============================================================================

def _add_footer_band(self, slide, content, slide_number, client_name="", month=""):
    """Two-line footer at the bottom of content slides."""
    if content.slide_type in ("title", "section"):
        return

    slide_w = Inches(13.33)
    slide_h = Inches(7.5)
    footer_left = Inches(0.5)
    footer_width = slide_w - Inches(1.0)

    # Line 1: source/methodology
    source_text = ""
    if content.notes_text:
        first_line = content.notes_text.split("\n")[0]
        if not first_line.startswith("KEY FINDING:"):
            source_text = first_line[:120]
    if not source_text:
        source_text = "Source: Client ODD file | Analysis by Velocity Solutions"

    line1_top = slide_h - Inches(0.5)
    tb = slide.shapes.add_textbox(footer_left, line1_top, footer_width, Inches(0.2))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = source_text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
    p.alignment = PP_ALIGN.LEFT

    # Line 2: Client | Month | Slide n | STRICTLY CONFIDENTIAL
    month_display = month
    if month and "." in month:
        try:
            import calendar as _cal
            yyyy, mm = month.split(".")[:2]
            month_display = f"{_cal.month_name[int(mm)]} {yyyy}"
        except (ValueError, IndexError):
            pass

    parts = []
    if client_name:
        parts.append(client_name)
    if month_display:
        parts.append(month_display)
    parts.append(f"Slide {slide_number}")
    parts.append("STRICTLY CONFIDENTIAL")
    footer_text = "  |  ".join(parts)

    line2_top = slide_h - Inches(0.28)
    tb = slide.shapes.add_textbox(footer_left, line2_top, footer_width, Inches(0.2))
    tb.text_frame.word_wrap = False
    p = tb.text_frame.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.LEFT


# =============================================================================
# 18.3 — Section divider
# REPLACE the existing DeckBuilder._build_section_slide (line ~428).
# SLIDE_DESIGN.md §9.
#
# Expects content.title in one of:
#   "02\nDebit Card Take Rate\nLead-in sentence"   (full 3-line)
#   "Debit Card Take Rate\nLead-in sentence"       (no number)
#   "Debit Card Take Rate"                         (title only)
# =============================================================================

def _build_section_slide(self, slide, content):
    # FIX (review bug 9): collect placeholders FIRST, then remove. Mutating
    # slide.placeholders during iteration can skip elements.
    placeholders_to_remove = [ph for ph in slide.placeholders]
    for ph in placeholders_to_remove:
        try:
            ph.element.getparent().remove(ph.element)
        except Exception:
            pass

    # Full-bleed navy background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x36, 0x5D)

    lines = (content.title or "").split("\n")
    section_number = None
    section_title = lines[0] if lines else ""
    lead_in = None

    if len(lines) >= 3:
        first = lines[0].strip()
        # Heuristic: a section number is 1-3 chars, all digits or "0X"
        if len(first) <= 3 and (first.isdigit() or (first.startswith("0") and first[1:].isdigit())):
            section_number = first
            section_title = lines[1]
            lead_in = "\n".join(lines[2:])
        else:
            section_title = lines[0]
            lead_in = "\n".join(lines[1:])
    elif len(lines) == 2:
        section_title = lines[0]
        lead_in = lines[1]

    y_cursor = Inches(2.5)

    if section_number:
        tb = slide.shapes.add_textbox(Inches(0.86), Inches(2.2), Inches(2.0), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = section_number
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0D, 0x94, 0x88)
        p.alignment = PP_ALIGN.LEFT
        y_cursor = Inches(3.0)

    tb = slide.shapes.add_textbox(Inches(0.86), y_cursor, Inches(11.0), Inches(1.0))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT

    if lead_in:
        tb = slide.shapes.add_textbox(
            Inches(0.86), y_cursor + Inches(1.1), Inches(11.0), Inches(1.0)
        )
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = lead_in
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.LEFT


# =============================================================================
# Modified _add_slide — REPLACE existing DeckBuilder._add_slide (line ~196).
# Adds slide_number, client_name, month kwargs so callout + footer can render.
# =============================================================================

def _add_slide(self, content, slide_number=0, client_name="", month=""):
    slide = self.prs.slides.add_slide(self.prs.slide_layouts[content.layout_index])

    builders = {
        "title": self._build_title_slide,
        "section": self._build_section_slide,
        "screenshot": self._build_screenshot_slide,
        "screenshot_kpi": self._build_screenshot_kpi_slide,
        "multi_screenshot": self._build_multi_screenshot_slide,
        "summary": self._build_summary_slide,
        "blank": self._build_blank_slide,
        "mailer_summary": self._build_mailer_summary_slide,
        "kpi_dashboard": self._build_kpi_dashboard_slide,
        "chart_narrative": self._build_chart_narrative_slide,
        "kpi_hero": self._build_kpi_hero_slide,
    }

    builder = builders.get(content.slide_type)
    if builder:
        builder(slide, content)
    else:
        logger.warning("Unknown slide_type: {t}", t=content.slide_type)

    # 18.1 — callout on chart slides (screenshot_kpi already shows KPIs natively)
    if content.slide_type in ("screenshot", "multi_screenshot") and content.kpis:
        self._add_callout_box(slide, content.kpis)

    # 18.2 — footer band on every content slide
    if slide_number > 0:
        self._add_footer_band(slide, content, slide_number, client_name, month)

    if content.notes_text:
        try:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                tf.text = content.notes_text
        except Exception:
            pass


# =============================================================================
# Modified build() — REPLACE existing DeckBuilder.build (line ~158).
# Threads client_name + month through to _add_slide for footer rendering.
# =============================================================================

def build(self, slides, output_path, client_name="", month=""):
    self.prs = Presentation(self.template_path)

    while len(self.prs.slides) > 0:
        rId = self.prs.slides._sldIdLst[0].get(qn("r:id"))
        self.prs.part.drop_rel(rId)
        self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    n_layouts = len(self.prs.slide_layouts)

    for i, slide_content in enumerate(slides):
        if slide_content.layout_index >= n_layouts:
            logger.warning(
                "Slide {i} '{title}' has layout_index={idx} but template only has {n} layouts, using 0",
                i=i, title=slide_content.title[:40], idx=slide_content.layout_index, n=n_layouts,
            )
            slide_content.layout_index = 0
        try:
            self._add_slide(
                slide_content,
                slide_number=i + 1,
                client_name=client_name,
                month=month,
            )
        except Exception as exc:
            logger.error(
                "Slide {i} '{title}' (type={t}, layout={l}) failed: {err}",
                i=i, title=slide_content.title[:40], t=slide_content.slide_type,
                l=slide_content.layout_index, err=exc,
            )

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    self.prs.save(output_path)
    return output_path


# =============================================================================
# build_deck() caller update (line ~1730 area)
# =============================================================================
# FIND:
#     result = builder.build(all_slides, str(pptx_path))
# REPLACE WITH:
#     result = builder.build(
#         all_slides,
#         str(pptx_path),
#         client_name=ctx.client.client_name,
#         month=ctx.client.month,
#     )


# =============================================================================
# _section_divider() helper update (inside build_deck, line ~1824)
# =============================================================================
# REVIEW FIX: the existing helper today is:
#
#   def _section_divider(title, subtitle=None, layout_index=LAYOUT_SECTION,
#                        slide_type="section"):
#       full_title = f"{title}\n{subtitle}" if subtitle else title
#       return SlideContent(slide_type=slide_type, title=full_title,
#                           layout_index=layout_index)
#
# Note: existing call sites pass layout_index=_DEFAULT_DIVIDER_LAYOUT
# (= LAYOUT_TITLE) for main-body section dividers. _build_section_slide
# now paints the background programmatically, so the layout choice no
# longer matters for the visual — we'll keep LAYOUT_SECTION as the
# argument for clarity going forward.
#
# REPLACE WITH:
#
#   def _section_divider(section_key, layout_index=LAYOUT_SECTION,
#                        slide_type="section"):
#       num = _SECTION_NUMBERS.get(section_key, "")
#       label = _SECTION_LABELS.get(section_key, section_key.title())
#       lead = _SECTION_LEAD_INS.get(section_key, "")
#       parts = [p for p in (num, label, lead) if p]
#       return SlideContent(
#           slide_type=slide_type,
#           title="\n".join(parts),
#           layout_index=layout_index,
#       )
#
# AND update its call sites in build_deck. The current loop is roughly:
#
#       for section_key in SECTION_ORDER:
#           slides = _section_main.get(section_key, [])
#           if not slides:
#               continue
#           label = _SECTION_LABELS.get(section_key, section_key.title())
#           analysis_slides.append(
#               _section_divider(label, subtitle=section_subtitle, layout_index=_DEFAULT_DIVIDER_LAYOUT)
#           )
#           analysis_slides.extend(slides)
#
# REPLACE WITH:
#
#       for section_key in SECTION_ORDER:
#           slides = _section_main.get(section_key, [])
#           if not slides:
#               continue
#           analysis_slides.append(_section_divider(section_key))
#           analysis_slides.extend(slides)
#
# The summary divider at the bottom of build_deck remains as-is:
#
#       _section_divider("Summary & Key Takeaways", layout_index=LAYOUT_SECTION, ...)
#
# That call passes a literal title string, not a section_key — handle that
# in _section_divider by checking if the arg looks like a known section key:
#
#   def _section_divider(section_key_or_title, layout_index=LAYOUT_SECTION,
#                        slide_type="section"):
#       if section_key_or_title in _SECTION_LABELS:
#           num = _SECTION_NUMBERS.get(section_key_or_title, "")
#           label = _SECTION_LABELS[section_key_or_title]
#           lead = _SECTION_LEAD_INS.get(section_key_or_title, "")
#           parts = [p for p in (num, label, lead) if p]
#           title = "\n".join(parts)
#       else:
#           title = section_key_or_title  # literal title (e.g. "Summary & Key Takeaways")
#       return SlideContent(slide_type=slide_type, title=title, layout_index=layout_index)
