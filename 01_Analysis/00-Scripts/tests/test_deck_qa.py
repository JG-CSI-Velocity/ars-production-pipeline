"""Tests for the deck QA gate.

Decks are synthesized in-memory so no client data is committed. Each test pins one
defect class that shipped in deck 1759 and must never regress silently again.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from ars_analysis.output import deck_qa as qa


def _blank_deck():
    prs = Presentation()
    prs._blank_layout = prs.slide_layouts[6]  # "Blank" — not a divider hint
    return prs


def _add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(slide, text, left=1.0, top=1.0, width=6.0, height=1.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.text = text
    return box


def _save(prs, tmp_path, name="deck.pptx"):
    path = tmp_path / name
    prs.save(str(path))
    return path


def test_clean_deck_passes(tmp_path):
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Quarterly Results", top=0.5, height=0.8)
    _add_text(slide, "Revenue grew steadily across all segments.", top=2.0, width=8.0, height=1.5)
    report = qa.audit_deck(_save(prs, tmp_path))
    assert report["passed"]
    assert report["counts"]["CRITICAL"] == 0
    assert report["counts"]["MAJOR"] == 0


def test_leaked_template_token_is_critical(tmp_path):
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Headline")
    _add_text(slide, "{overall_rate:.1f}% across {total_mailed:,} mailed", top=3.0, width=8.0)
    report = qa.audit_deck(_save(prs, tmp_path))
    codes = {f["code"] for f in report["findings"]}
    assert "leaked_token" in codes
    assert report["counts"]["CRITICAL"] >= 1
    assert not report["passed"]


def test_slide_count_explosion_flagged(tmp_path):
    prs = _blank_deck()
    for _ in range(qa.MAX_SLIDES + 1):
        _add_slide(prs)
    report = qa.audit_deck(_save(prs, tmp_path))
    assert any(f["code"] == "slide_count" for f in report["findings"])


def test_text_overflow_flagged(tmp_path):
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Title", top=0.5)
    _add_text(slide, "x" * 120, top=2.0, width=1.4, height=0.4)  # the 1.4in mailer stat box
    report = qa.audit_deck(_save(prs, tmp_path))
    assert any(f["code"] == "text_overflow" for f in report["findings"])


def test_empty_body_flagged(tmp_path):
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Deposit Trends", top=0.5)  # title only, no body, no chart
    report = qa.audit_deck(_save(prs, tmp_path))
    assert any(f["code"] == "empty_body" for f in report["findings"])


def test_operator_filled_slides_not_flagged(tmp_path):
    """Agenda / Exec Summary / Monthly Revenue / ARS Lift are blank by design."""
    prs = _blank_deck()
    for title in ("Agenda", "Executive Summary",
                  "Monthly Revenue – Last 12 Months", "ARS Lift Matrix"):
        _add_text(_add_slide(prs), title, top=0.5)
    report = qa.audit_deck(_save(prs, tmp_path))
    assert not any(f["code"] == "empty_body" for f in report["findings"])


def test_zero_denominator_is_critical(tmp_path):
    """The A9.6 '0 / 0' bar and any 'n / 0' stat must fail the gate (#217)."""
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Attrition Comparison", top=0.5)
    _add_text(slide, "Closure rate 0.0% (0 / 0)", top=2.0, width=8.0)
    report = qa.audit_deck(_save(prs, tmp_path))
    codes = {f["code"] for f in report["findings"]}
    assert "zero_denominator" in codes
    assert not report["passed"]


def test_implausible_split_flagged(tmp_path):
    """One arm 0.0% vs another >=20% on a stat-pair slide (#217, the A9.9 debit split)."""
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Debit Card Impact on Retention", top=0.5)
    _add_text(slide, "With Debit Card: 0.0% (4 / 14,642)", top=2.0, width=8.0)
    _add_text(slide, "Without Debit Card: 29.2% (544 / 1,861)", top=3.0, width=8.0)
    report = qa.audit_deck(_save(prs, tmp_path))
    assert any(f["code"] == "implausible_split" for f in report["findings"])


def test_plain_percentages_not_flagged_as_split(tmp_path):
    """A 0.0% growth figure next to an unrelated large rate must NOT trip the
    split check -- only pct(n/d) stat pairs are compared."""
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Deposit Overview", top=0.5)
    _add_text(slide, "Growth was 0.0% this month; take rate stands at 35.2%.",
              top=2.0, width=8.0)
    report = qa.audit_deck(_save(prs, tmp_path))
    assert not any(f["code"] == "implausible_split" for f in report["findings"])


def test_healthy_stat_pair_not_flagged(tmp_path):
    prs = _blank_deck()
    slide = _add_slide(prs)
    _add_text(slide, "Debit Card Impact on Retention", top=0.5)
    _add_text(slide, "With Debit Card: 12.1% (500 / 4,132)", top=2.0, width=8.0)
    _add_text(slide, "Without Debit Card: 18.9% (351 / 1,861)", top=3.0, width=8.0)
    report = qa.audit_deck(_save(prs, tmp_path))
    assert not any(f["code"] in ("implausible_split", "zero_denominator")
                   for f in report["findings"])
