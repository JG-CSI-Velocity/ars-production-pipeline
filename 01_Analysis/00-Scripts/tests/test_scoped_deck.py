"""build_scoped_deck must produce a one-section deck (title + divider + that
section's slides only) in a flat modules/ folder, using the real template, and
apply the same consolidation the full report does (merge paired slides, split an
appendix). This is the deliverable for the per-module closed loop."""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from pptx import Presentation  # noqa: E402

from ars_analysis.analytics.base import AnalysisResult  # noqa: E402
from ars_analysis.analytics.section_registry import get_section  # noqa: E402
from ars_analysis.output import deck_builder  # noqa: E402
from ars_analysis.pipeline.context import (  # noqa: E402
    ClientInfo,
    OutputPaths,
    PipelineContext,
)


def _png(path: Path) -> Path:
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.bar([0, 1, 2], [3, 1, 2])
    fig.savefig(path, dpi=100)
    plt.close(fig)
    return path


def _ctx(tmp_path: Path, slides) -> PipelineContext:
    ctx = PipelineContext(
        client=ClientInfo(client_id="1776", client_name="CoastHills CU",
                          month="2026.06", assigned_csm="JamesG"),
        paths=OutputPaths.from_dir(tmp_path),
    )
    ctx.all_slides = slides
    return ctx


def test_scoped_txn_deck_contains_only_its_section(tmp_path):
    slides = [
        AnalysisResult(slide_id="TXN-MERCH-01", title="Top Merchants",
                       chart_path=_png(tmp_path / "m1.png"), slide_type="screenshot"),
        AnalysisResult(slide_id="TXN-MERCH-02", title="Merchant Concentration",
                       chart_path=_png(tmp_path / "m2.png"), slide_type="screenshot"),
        # A different section's slide must NOT appear in the merchant deck.
        AnalysisResult(slide_id="TXN-COMP-01", title="Competitor Detection",
                       chart_path=_png(tmp_path / "c1.png"), slide_type="screenshot"),
    ]
    ctx = _ctx(tmp_path, slides)
    section = get_section("txn.merchant")

    out = deck_builder.build_scoped_deck(ctx, section)

    assert out is not None and out.exists()
    # Flat modules/ folder (no per-module subfolder) + module-scoped filename.
    assert out.parent.name == "modules"
    assert out.name == "1776_2026.06_txn_merchant_deck.pptx"
    # title + divider + 2 merchant slides == 4 (the competition slide excluded).
    prs = Presentation(str(out))
    assert len(prs.slides) == 4


def _slide_texts(prs) -> list[str]:
    out = []
    for slide in prs.slides:
        parts = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame]
        out.append(" ".join(parts))
    return out


def test_scoped_dctr_deck_merges_funnel_and_splits_appendix(tmp_path):
    """A DCTR module deck must match the full report: the funnel pair
    (A7.7 + A7.8) merges into one 2x1 slide and a DCTR_APPENDIX_IDS slide (A7.5)
    drops behind an Appendix divider -- not shipped as bare single slides."""
    slides = [
        AnalysisResult(slide_id="A7.1", title="DCTR Overview",
                       chart_path=_png(tmp_path / "d1.png"), slide_type="screenshot"),
        AnalysisResult(slide_id="A7.7", title="Funnel Historical",
                       chart_path=_png(tmp_path / "d7.png"), slide_type="screenshot"),
        AnalysisResult(slide_id="A7.8", title="Funnel TTM",
                       chart_path=_png(tmp_path / "d8.png"), slide_type="screenshot"),
        AnalysisResult(slide_id="A7.5", title="DCTR Detail",
                       chart_path=_png(tmp_path / "d5.png"), slide_type="screenshot"),
    ]
    ctx = _ctx(tmp_path, slides)
    out = deck_builder.build_scoped_deck(ctx, get_section("ars.dctr"))

    assert out is not None and out.exists()
    assert out.parent.name == "modules"
    assert out.name == "1776_2026.06_ars_dctr_deck.pptx"

    prs = Presentation(str(out))
    # title, section divider, A7.1, merged funnel, Appendix divider, A7.5 == 6.
    # (Unmerged would be 7; no appendix split would drop the divider.)
    assert len(prs.slides) == 6
    texts = _slide_texts(prs)
    assert any("Appendix" in t for t in texts), "appendix divider missing"
    assert any("Funnel" in t for t in texts), "merged funnel slide missing"


def test_scoped_deck_none_when_section_absent(tmp_path):
    slides = [
        AnalysisResult(slide_id="TXN-MERCH-01", title="Top Merchants",
                       chart_path=_png(tmp_path / "m1.png"), slide_type="screenshot"),
    ]
    ctx = _ctx(tmp_path, slides)
    # No ICS_cohort slides present -> nothing to build.
    out = deck_builder.build_scoped_deck(ctx, get_section("txn.ICS_cohort"))
    assert out is None
