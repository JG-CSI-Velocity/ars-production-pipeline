"""The deck QA gate must actually run on a finished build.

deck_qa.audit_deck shipped with full coverage but was only ever called from
tests -- which is how deck 1759 shipped with leaked template tokens. These tests
lock the wiring in build_deck: a quality report is written next to every deck,
and a failing deck escalates to a manifest anomaly flag.
"""

from __future__ import annotations

from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches

from ars_analysis.output.deck_builder import _run_deck_qa
from ars_analysis.pipeline.manifest import FlagLevel


class _StubManifest:
    """Minimal manifest exposing the flag() API _run_deck_qa calls."""

    def __init__(self):
        self.flags = []

    def flag(self, level: FlagLevel, message: str) -> None:
        self.flags.append((level, message))


def _deck_with(tmp_path, name, body_text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    box.text_frame.text = "Headline"
    box2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    box2.text_frame.text = body_text
    path = tmp_path / name
    prs.save(str(path))
    return path


def test_qa_writes_report_and_flags_on_leak(tmp_path):
    deck = _deck_with(tmp_path, "1759_2026.06_ars_deck.pptx",
                      "{overall_rate:.1f}% across {total_mailed:,} mailed")
    manifest = _StubManifest()
    ctx = SimpleNamespace(export_log=[], manifest=manifest)

    report = _run_deck_qa(deck, ctx)

    assert report is not None and not report["passed"]
    report_path = deck.with_name(f"{deck.stem}_quality_report.txt")
    assert report_path.exists()
    assert str(report_path) in ctx.export_log
    # A leaked token is CRITICAL -> manifest gets an ERROR-level flag.
    assert manifest.flags and manifest.flags[0][0] == FlagLevel.ERROR


def test_qa_clean_deck_no_flag(tmp_path):
    deck = _deck_with(tmp_path, "1759_2026.06_ars_deck.pptx",
                      "Revenue grew steadily across every segment this period.")
    manifest = _StubManifest()
    ctx = SimpleNamespace(export_log=[], manifest=manifest)

    report = _run_deck_qa(deck, ctx)

    assert report is not None and report["passed"]
    assert manifest.flags == []


def test_qa_survives_missing_manifest(tmp_path):
    """No manifest on ctx must not raise -- QA is best-effort, never build-breaking."""
    deck = _deck_with(tmp_path, "1759_2026.06_ars_deck.pptx",
                      "{leaked_token} should still produce a report")
    ctx = SimpleNamespace(export_log=[], manifest=None)

    report = _run_deck_qa(deck, ctx)

    assert report is not None
    assert deck.with_name(f"{deck.stem}_quality_report.txt").exists()
