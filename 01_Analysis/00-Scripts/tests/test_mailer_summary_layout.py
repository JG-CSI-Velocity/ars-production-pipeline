"""Locks the mailer-summary slide layout against the 1759 overlap regressions.

Builds a real mailer-summary slide from the CSI template with content that mixes
full-sentence insights and value|description stats (the exact shape that broke),
then runs the deck QA gate. Two regressions must stay fixed:
  - full sentences must not be crammed into the 1.4in value box (text overflow)
  - the layout's empty TITLE placeholder must not linger (orphan overlap)
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from ars_analysis.output import deck_qa
from ars_analysis.output.deck_builder import DeckBuilder, SlideContent, LAYOUT_PICTURE

_TEMPLATE = (
    Path(__file__).resolve().parents[1] / "output" / "template" / "2025-CSI-PPT-Template.pptx"
)

_BULLETS = [
    "TH-25 led with a 31.5% response rate (132 of 419 mailed)",
    "TH-25 responded at 6.5x the rate of NU 5+ (4.9%) -- higher tiers drive response",
    "37% of responders were new to the program this wave",
    "23%|of Responders were accounts opened fewer than 12 months ago",
    "24%|of Responders aged 60+",
    "26%|of Responders opted into Reg E",
    "37%|First-time responders",
    "40%|of repeat responders moved up the ladder",
]


@pytest.fixture
def mailer_deck(tmp_path):
    if not _TEMPLATE.exists():
        pytest.skip(f"template not present: {_TEMPLATE}")
    builder = DeckBuilder(str(_TEMPLATE))
    builder.prs = Presentation(str(_TEMPLATE))
    while len(builder.prs.slides) > 0:
        rid = builder.prs.slides._sldIdLst[0].get(qn("r:id"))
        builder.prs.part.drop_rel(rid)
        builder.prs.slides._sldIdLst.remove(builder.prs.slides._sldIdLst[0])
    content = SlideContent(
        slide_type="mailer_summary",
        title="Apr26 mailer: 5,564 mailed, 13.3% response",
        layout_index=LAYOUT_PICTURE,
        title_color="#1B365D",
        kpis={"Mailed": "5,564", "Responded": "741", "Rate": "13.3%"},
        bullets=_BULLETS,
        images=[],
    )
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[LAYOUT_PICTURE])
    builder._build_mailer_summary_slide(slide, content)
    out = tmp_path / "mailer.pptx"
    builder.prs.save(str(out))
    return out


def test_no_text_overflow(mailer_deck):
    report = deck_qa.audit_deck(mailer_deck)
    overflow = [f for f in report["findings"] if f["code"] == "text_overflow"]
    assert overflow == [], overflow


def test_no_leftover_placeholders(mailer_deck):
    prs = Presentation(str(mailer_deck))
    assert len(list(prs.slides[0].placeholders)) == 0


def test_mailer_performance_ancillary_deck_is_written(tmp_path):
    """Older waves produce a separate 'Mailer Performance' .pptx next to the main deck."""
    if not _TEMPLATE.exists():
        pytest.skip(f"template not present: {_TEMPLATE}")
    from types import SimpleNamespace
    from ars_analysis.output.deck_builder import _build_mailer_performance_deck

    ctx = SimpleNamespace(
        client=SimpleNamespace(client_id="1759", month="2026.06"),
        export_log=[],
    )
    older = [
        SlideContent(slide_type="mailer_summary", title="Apr25 mailer: 5,160 mailed, 14.8% response",
                     layout_index=LAYOUT_PICTURE, bullets=["TH-25 led", "16%|opted in"], images=[]),
        SlideContent(slide_type="title", title="Apr25 combo", layout_index=LAYOUT_PICTURE),
    ]
    main_output = tmp_path / "1759_2026.06_ars_deck.pptx"
    out = _build_mailer_performance_deck(
        ctx, older, "First Central CU", "First Central CU | June 2026",
        str(_TEMPLATE), main_output,
    )
    assert out is not None and out.exists()
    assert out.name == "1759_2026.06_Mailer_Performance.pptx"
    # cover + divider + the 2 older-wave slides
    assert len(Presentation(str(out)).slides) == 4


def test_mailer_performance_deck_noop_when_no_older_waves(tmp_path):
    from types import SimpleNamespace
    from ars_analysis.output.deck_builder import _build_mailer_performance_deck
    ctx = SimpleNamespace(client=SimpleNamespace(client_id="1759", month="2026.06"), export_log=[])
    out = _build_mailer_performance_deck(
        ctx, [], "First Central CU", "sub", "ignored.pptx", tmp_path / "main.pptx",
    )
    assert out is None
