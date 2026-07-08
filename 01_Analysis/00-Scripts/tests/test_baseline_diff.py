"""diff_decks must report two decks identical when their per-slide text +
picture counts match, and flag real differences -- the content check behind the
one-command baseline-diff the M: box runs to confirm the refactor is output-safe.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "tools"))
import baseline_diff  # noqa: E402


def _deck(path: Path, titles, pics_per_slide=0, tiny_png=None):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for t in titles:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tb.text_frame.text = t
        for _ in range(pics_per_slide):
            if tiny_png:
                s.shapes.add_picture(str(tiny_png), Inches(1), Inches(3),
                                     Inches(1), Inches(1))
    prs.save(str(path))
    return path


def _png(path: Path):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=(2, 2))
    ax.plot([0, 1], [0, 1])
    fig.savefig(path)
    plt.close(fig)
    return path


def test_identical_decks_report_identical(tmp_path):
    a = _deck(tmp_path / "a.pptx", ["Merchant Analysis", "Top Merchants"])
    b = _deck(tmp_path / "b.pptx", ["Merchant Analysis", "Top Merchants"])
    d = baseline_diff.diff_decks(a, b)
    assert d.identical
    assert d.slide_count_a == d.slide_count_b == 2


def test_text_difference_is_flagged(tmp_path):
    a = _deck(tmp_path / "a.pptx", ["Merchant Analysis", "Top Merchants"])
    b = _deck(tmp_path / "b.pptx", ["Merchant Analysis", "DIFFERENT"])
    d = baseline_diff.diff_decks(a, b)
    assert not d.identical
    assert any("slide 2" in line and "text" in line for line in d.per_slide)


def test_slide_count_and_picture_diffs_are_flagged(tmp_path):
    png = _png(tmp_path / "p.png")
    a = _deck(tmp_path / "a.pptx", ["S1", "S2", "S3"], pics_per_slide=1, tiny_png=png)
    b = _deck(tmp_path / "b.pptx", ["S1", "S2"], pics_per_slide=0)
    d = baseline_diff.diff_decks(a, b)
    assert not d.identical
    assert d.slide_count_a == 3 and d.slide_count_b == 2
    assert any("present in current only" in line for line in d.per_slide)
    assert any("picture count" in line for line in d.per_slide)
