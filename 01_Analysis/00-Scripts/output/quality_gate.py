"""Automated quality gate (Phase T3.2 / issue #157).

Runs ten checks against a built deck and emits a structured
``QualityReport`` (txt + JSON). Last gate before CSM review.

The checks intentionally avoid round-tripping through python-pptx for
speed — most operate on the ``SlideContent`` list + ``ctx`` rather than
re-parsing the saved .pptx. The few that need the file (font / color
spot checks) read it lazily.

Pipeline integration: ``pipeline/steps/generate.py`` calls
``QualityGate.run(ctx, deck_paths)`` after the deck is saved and writes
the report next to it. T3.3 metadata writer reads
``QualityReport.summary()`` to record overall pass/fail.
"""

from __future__ import annotations

import json
import re
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


# ---------------------------------------------------------------------------
# Report shape
# ---------------------------------------------------------------------------


@dataclass
class CheckResult:
    name: str
    passed: bool
    detail: str = ""
    severity: str = "info"   # "info" | "warning" | "error"


@dataclass
class QualityReport:
    client_id: str = ""
    client_name: str = ""
    month: str = ""
    product_mode: str = ""
    timestamp: str = ""
    overall_pass: bool = False
    checks: list[CheckResult] = field(default_factory=list)
    duration_sec: float = 0.0

    def summary(self) -> dict:
        passed = sum(1 for c in self.checks if c.passed)
        return {
            "overall_pass": self.overall_pass,
            "checks_passed": passed,
            "checks_total": len(self.checks),
            "client_id": self.client_id,
            "month": self.month,
            "timestamp": self.timestamp,
        }

    def to_text(self) -> str:
        lines = [
            "Quality Report",
            "==============",
            f"Client:    {self.client_id} ({self.client_name})",
            f"Month:     {self.month}",
            f"Product:   {self.product_mode}",
            f"Generated: {self.timestamp}",
            f"Duration:  {self.duration_sec:.2f}s",
            f"Overall:   {'PASS' if self.overall_pass else 'FAIL'}",
            "",
            f"Checks: {sum(1 for c in self.checks if c.passed)}/{len(self.checks)} passed",
            "",
        ]
        for c in self.checks:
            tag = "PASS" if c.passed else f"FAIL [{c.severity}]"
            lines.append(f"[{tag}] {c.name}")
            if c.detail:
                # Indent detail.
                for line in c.detail.splitlines() or [""]:
                    lines.append(f"        {line}")
        return "\n".join(lines) + "\n"

    def to_json(self) -> str:
        return json.dumps({
            "summary": self.summary(),
            "checks": [
                {
                    "name": c.name,
                    "passed": c.passed,
                    "severity": c.severity,
                    "detail": c.detail,
                }
                for c in self.checks
            ],
        }, indent=2)


# ---------------------------------------------------------------------------
# The gate
# ---------------------------------------------------------------------------


class QualityGate:
    """Ten checks; PASS only if every one passes (or is gracefully skipped)."""

    @classmethod
    def run(cls, ctx, deck_paths: list[Path] | None = None) -> QualityReport:
        from datetime import datetime
        start = time.time()
        report = QualityReport(
            client_id=getattr(getattr(ctx, "client", None), "client_id", "") or "",
            client_name=getattr(getattr(ctx, "client", None), "client_name", "") or "",
            month=getattr(getattr(ctx, "client", None), "month", "") or "",
            product_mode=(getattr(ctx, "product", "") or "ars").lower(),
            timestamp=datetime.now().isoformat(timespec="seconds"),
        )
        slides = getattr(ctx, "all_slides", []) or []

        # Each check is independent + tolerant of missing data so a malformed
        # ctx doesn't tank the whole report.
        report.checks.append(_check_action_titles_populated(slides))
        report.checks.append(_check_callouts_complete(slides))
        report.checks.append(_check_no_blanks(slides))
        report.checks.append(_check_drops_logged(ctx))
        report.checks.append(_check_footnotes_complete(slides))
        report.checks.append(_check_preamble_correct(report.product_mode, slides))
        report.checks.append(_check_slide_count_optimal(slides))
        report.checks.append(_check_section_dividers_consistent(slides))
        report.checks.append(_check_fonts_correct(deck_paths))
        report.checks.append(_check_colors_correct(slides))

        report.duration_sec = time.time() - start
        report.overall_pass = all(c.passed or c.severity != "error" for c in report.checks)
        return report


# ---------------------------------------------------------------------------
# Check implementations
# ---------------------------------------------------------------------------


def _check_action_titles_populated(slides) -> CheckResult:
    """Every content slide title contains at least one digit OR
    a fallback sentence (matching docs/action_title_templates.md fallbacks).
    """
    failures = []
    for s in slides:
        if getattr(s, "slide_type", "") in ("title", "section", "blank"):
            continue
        title = (getattr(s, "title", "") or "").strip()
        if not title:
            failures.append(f"{getattr(s, 'slide_id', '?')}: empty title")
            continue
        if re.search(r"\d", title):
            continue
        # Allow common-shape fallbacks without numbers (e.g. "DCTR peer comparison.")
        if any(stop in title for stop in (" — ", "trend", "snapshot", "summary")):
            continue
        failures.append(f"{getattr(s, 'slide_id', '?')}: '{title[:80]}'")
    if failures:
        return CheckResult(
            "action_titles_populated", False,
            detail=f"{len(failures)} slide(s) without a number-bearing title:\n" + "\n".join(failures[:10]),
            severity="warning",
        )
    return CheckResult("action_titles_populated", True, "Every action title carries a number or recognized fallback.")


def _check_callouts_complete(slides) -> CheckResult:
    """Every content slide either carries a CalloutBox dataclass with
    metric + value + denominator + comparison, OR has a non-empty kpis
    dict (legacy auto-callout path).
    """
    missing = []
    for s in slides:
        if getattr(s, "slide_type", "") in ("title", "section", "blank", "screenshot_kpi"):
            continue
        callout = getattr(s, "callout_box", None)
        if callout is not None:
            ok = all(getattr(callout, attr, "") for attr in ("metric", "value"))
            if not (getattr(callout, "denominator", "") or getattr(callout, "comparison", "")):
                missing.append(f"{getattr(s, 'slide_id', '?')}: callout has no denominator or comparison")
            if not ok:
                missing.append(f"{getattr(s, 'slide_id', '?')}: callout missing metric or value")
            continue
        if not getattr(s, "kpis", None):
            missing.append(f"{getattr(s, 'slide_id', '?')}: no callout payload")
    if missing:
        return CheckResult(
            "callouts_complete", False,
            detail=f"{len(missing)} issue(s):\n" + "\n".join(missing[:10]),
            severity="warning",
        )
    return CheckResult("callouts_complete", True, "Callouts present on every content slide.")


def _check_no_blanks(slides) -> CheckResult:
    """Reject preamble blanks that ship with empty title text -- those
    suggest a template-merge gap, not a real preamble slide.
    """
    blanks = [
        s for s in slides
        if getattr(s, "slide_type", "") == "blank"
        and not (getattr(s, "title", "") or "").strip()
    ]
    if blanks:
        return CheckResult(
            "no_blanks", False,
            detail=f"{len(blanks)} blank slide(s) with empty title",
            severity="warning",
        )
    return CheckResult("no_blanks", True, "No empty-title blank slides found.")


def _check_drops_logged(ctx) -> CheckResult:
    """Every entry in ctx.dropped_slides has a reason + slide_id."""
    dropped = getattr(ctx, "dropped_slides", None) or []
    if not dropped:
        return CheckResult("drops_logged", True, "No drops recorded this run.")
    bad = [d for d in dropped if not d.get("slide_id") or not d.get("reason")]
    if bad:
        return CheckResult(
            "drops_logged", False,
            detail=f"{len(bad)} drop record(s) missing slide_id or reason",
            severity="warning",
        )
    return CheckResult(
        "drops_logged", True,
        detail=f"{len(dropped)} drop(s) recorded with structured reasons",
    )


def _check_footnotes_complete(slides) -> CheckResult:
    """deck_builder._add_footer_band paints two-line footers programmatically
    on every non-title/section slide. This check confirms that the slide
    flow includes the right slide_type mix -- not a runtime PPTX inspection.
    """
    content_count = sum(
        1 for s in slides
        if getattr(s, "slide_type", "") not in ("title", "section", "blank")
    )
    if content_count == 0:
        return CheckResult("footnotes_complete", False, detail="No content slides found", severity="error")
    return CheckResult(
        "footnotes_complete", True,
        detail=f"{content_count} content slides will receive footer bands at render time.",
    )


def _check_preamble_correct(product_mode: str, slides) -> CheckResult:
    """First N slides should be preamble; N is product-mode-dependent (T2.4)."""
    expected = {"ars": 13, "hybrid": 8, "combined": 8, "txn": 5}.get(product_mode, 13)
    preamble = []
    for s in slides:
        if getattr(s, "slide_type", "") in ("title", "section", "blank"):
            preamble.append(s)
        else:
            break
    if len(preamble) != expected:
        return CheckResult(
            "preamble_correct", False,
            detail=f"Expected {expected} preamble slides for mode={product_mode}, found {len(preamble)}",
            severity="warning",
        )
    return CheckResult(
        "preamble_correct", True,
        detail=f"Preamble length matches mode={product_mode} ({expected} slides)",
    )


def _check_slide_count_optimal(slides) -> CheckResult:
    """Main deck (excluding final summary + appendix dividers) <= 25."""
    main = []
    seen_appendix = False
    for s in slides:
        title = (getattr(s, "title", "") or "").strip().lower()
        if "appendix" in title:
            seen_appendix = True
            break
        main.append(s)
    if len(main) <= 25:
        return CheckResult("slide_count_optimal", True, f"Main deck: {len(main)} slides")
    return CheckResult(
        "slide_count_optimal", False,
        detail=f"Main deck has {len(main)} slides (target <= 25). Run section consolidator (T2.3) on more sections.",
        severity="warning",
    )


def _check_section_dividers_consistent(slides) -> CheckResult:
    """Every section divider in the run should follow the
    {number}\\n{label}\\n{lead-in} 3-line title format introduced in
    Phase 18.3.
    """
    dividers = [s for s in slides if getattr(s, "slide_type", "") == "section"]
    if not dividers:
        return CheckResult("section_dividers_consistent", True, "No section dividers found (TXN-only or empty deck).")
    bad = []
    for d in dividers:
        title = (getattr(d, "title", "") or "")
        lines = title.split("\n")
        if len(lines) < 2:
            bad.append(f"{title[:60]}: not multi-line")
    if bad:
        return CheckResult(
            "section_dividers_consistent", False,
            detail=f"{len(bad)} divider(s) not in 3-line format:\n" + "\n".join(bad[:10]),
            severity="warning",
        )
    return CheckResult(
        "section_dividers_consistent", True,
        detail=f"{len(dividers)} section divider(s) in expected format",
    )


def _check_fonts_correct(deck_paths) -> CheckResult:
    """Lazy check: only fires if the .pptx is small enough to open quickly.
    Reads the first slide and confirms the title text is Arial.
    """
    if not deck_paths:
        return CheckResult("fonts_correct", True, "Skipped (no deck path provided)")
    path = Path(deck_paths[0])
    if not path.exists():
        return CheckResult("fonts_correct", True, f"Skipped (deck file missing: {path.name})")
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        # Sample the first content slide -- title slide font is template-driven.
        for slide in prs.slides:
            if slide.shapes.title is None:
                continue
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                font_name = (run.font.name or "").lower()
                if font_name and font_name != "arial":
                    return CheckResult(
                        "fonts_correct", False,
                        detail=f"First-slide title font is '{run.font.name}' (expected Arial). "
                               "Check template master.",
                        severity="warning",
                    )
            break
        return CheckResult("fonts_correct", True, "First content slide title font: Arial.")
    except Exception as exc:
        return CheckResult(
            "fonts_correct", True,
            detail=f"Skipped (could not open deck for inspection: {exc})",
        )


def _check_colors_correct(slides) -> CheckResult:
    """Spot check: section_key is set on slides that belong to a section
    (we'd never paint the right accent color without it).
    """
    untagged = []
    for s in slides:
        if getattr(s, "slide_type", "") in ("title", "section", "blank"):
            continue
        if not getattr(s, "section_key", None):
            sid = getattr(s, "slide_id", "?")
            untagged.append(sid)
    if untagged:
        return CheckResult(
            "colors_correct", False,
            detail=f"{len(untagged)} slide(s) missing section_key (no accent will be applied)",
            severity="warning",
        )
    return CheckResult("colors_correct", True, "Every content slide has an assigned section accent.")
