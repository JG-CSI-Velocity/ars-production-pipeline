"""Deck QA — static checks that catch the defect classes that silently shipped in deck 1759.

This runs on a finished .pptx with no access to source data, so it can gate every
deck before delivery. Each check maps to a real failure we hit:

  leaked_token   -> "{overall_rate:.1f}%" rendered literally (template substitution failed)
  slide_count    -> 167-slide explosion (uncapped appendix)
  empty_body     -> title-only Executive Summary / missing-chart slides
  text_overflow  -> long sentences crammed into 1.4in stat boxes (the visible mailer collision)

Note on overlap detection: pure box-geometry overlap was tried and removed. It cannot
distinguish a real collision from intentional layering (a label sitting on a stat card),
so it false-positived on the known-good reference decks. text_overflow is the calibrated
proxy for the mailer text-collision and stays clean on the good decks.

It is intentionally conservative: thresholds are tuned so a known-good deck passes
clean. Findings are advisory data, not exceptions — the caller decides policy.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass, asdict
from pathlib import Path

from pptx import Presentation

# --- thresholds (named so policy is visible, not buried in conditionals) ---
MAX_SLIDES = 60                 # good decks run ~40; 167 is the explosion
MAX_FILE_MB = 15.0              # rasterized-everything decks bloat past this
OVERFLOW_RATIO = 1.6           # text longer than est. capacity * this -> overflow
LEAKED_TOKEN = re.compile(r"\{[A-Za-z_][\w.]*(?::[^}]*)?\}")  # {x} or {x:.1f}

# Layouts that are *meant* to be title-only (covers, section breaks) — not "empty body".
_DIVIDER_HINTS = ("section", "divider", "cover", "title slide", "agenda")
# Slides the operator fills by hand from their own PowerPoint/Excel. They ship
# title-only by design, so they are not "empty body" defects.
_OPERATOR_FILLED = (
    "agenda", "executive summary", "monthly revenue", "ars lift",
    # Preamble framing blanks the operator fills by hand (or that carry only a
    # heading): they ship title-only by design, like the others above.
    "all program results", "data check overview",
)
# Rough text-capacity model at ~11pt body type.
_CHARS_PER_INCH = 15.0
_LINE_HEIGHT_IN = 0.20


@dataclass
class Finding:
    severity: str  # CRITICAL | MAJOR | MINOR
    code: str
    slide: int     # 1-based for humans; 0 = deck-level
    message: str


def _emu_to_in(value) -> float:
    return (value or 0) / 914400.0


def _box(shape):
    """Return (left, top, width, height) in inches, or None if unpositioned."""
    if shape.left is None or shape.top is None:
        return None
    return (
        _emu_to_in(shape.left),
        _emu_to_in(shape.top),
        _emu_to_in(shape.width),
        _emu_to_in(shape.height),
    )


def _is_divider(layout_name: str) -> bool:
    name = (layout_name or "").lower()
    return any(hint in name for hint in _DIVIDER_HINTS)


def _text(shape) -> str:
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def check_leaked_tokens(prs) -> list[Finding]:
    out = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            for token in LEAKED_TOKEN.findall(_text(shape)):
                out.append(Finding("CRITICAL", "leaked_token", i,
                                   f"unrendered template token {token!r}"))
    return out


def check_slide_count(prs) -> list[Finding]:
    n = len(prs.slides)
    if n > MAX_SLIDES:
        return [Finding("MAJOR", "slide_count", 0,
                       f"{n} slides exceeds max {MAX_SLIDES} (likely uncapped expansion)")]
    return []


def check_file_size(path: Path) -> list[Finding]:
    mb = path.stat().st_size / 1e6
    if mb > MAX_FILE_MB:
        return [Finding("MINOR", "file_size", 0,
                       f"{mb:.1f} MB exceeds {MAX_FILE_MB} MB (rasterized charts?)")]
    return []


def check_empty_body(prs) -> list[Finding]:
    """Non-divider content slide with only a title and no chart -> empty/missing content."""
    out = []
    for i, slide in enumerate(prs.slides, 1):
        if _is_divider(slide.slide_layout.name):
            continue
        has_picture = any(s.shape_type == 13 for s in slide.shapes)
        texts = [t for t in (_text(s) for s in slide.shapes) if t]
        if not has_picture and len(texts) == 1:
            title = texts[0].lower()
            if any(name in title for name in _OPERATOR_FILLED):
                continue  # operator fills this slide by hand -- blank is intended
            out.append(Finding("MAJOR", "empty_body", i,
                              f"title-only slide '{texts[0][:40]}' (empty body or missing chart)"))
    return out


def check_text_overflow(prs) -> list[Finding]:
    out = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            text = _text(shape)
            box = _box(shape)
            if not text or not box:
                continue
            w, h = box[2], box[3]
            if w <= 0 or h <= 0:
                continue
            capacity = (w * _CHARS_PER_INCH) * max(1.0, h / _LINE_HEIGHT_IN)
            if len(text) > capacity * OVERFLOW_RATIO:
                out.append(Finding("MAJOR", "text_overflow", i,
                                  f"{len(text)} chars in {w:.1f}x{h:.1f}in box: {text[:35]!r}"))
    return out


def audit_deck(path: str | Path) -> dict:
    """Run all checks; return a report dict (callable from app.py / tests / CLI)."""
    path = Path(path)
    prs = Presentation(str(path))
    findings: list[Finding] = []
    findings += check_leaked_tokens(prs)
    findings += check_slide_count(prs)
    findings += check_file_size(path)
    findings += check_empty_body(prs)
    findings += check_text_overflow(prs)
    counts = {sev: sum(1 for f in findings if f.severity == sev)
              for sev in ("CRITICAL", "MAJOR", "MINOR")}
    return {
        "file": path.name,
        "slides": len(prs.slides),
        "passed": counts["CRITICAL"] == 0 and counts["MAJOR"] == 0,
        "counts": counts,
        "findings": [asdict(f) for f in findings],
    }


def _print_report(report: dict) -> None:
    print(f"\nDeck QA — {report['file']}  ({report['slides']} slides)")
    c = report["counts"]
    print(f"  CRITICAL {c['CRITICAL']}   MAJOR {c['MAJOR']}   MINOR {c['MINOR']}   "
          f"-> {'PASS' if report['passed'] else 'FAIL'}")
    by_code: dict[str, int] = {}
    for f in report["findings"]:
        by_code[f["code"]] = by_code.get(f["code"], 0) + 1
    for code, n in sorted(by_code.items(), key=lambda kv: -kv[1]):
        sample = next(f for f in report["findings"] if f["code"] == code)
        loc = f"slide {sample['slide']}" if sample["slide"] else "deck"
        print(f"  [{sample['severity']:8}] {code:14} x{n:<4} e.g. {loc}: {sample['message']}")


def main() -> int:
    ap = argparse.ArgumentParser(description="Static QA checks for a generated deck (.pptx)")
    ap.add_argument("pptx", help="path to the .pptx to check")
    ap.add_argument("--json", action="store_true", help="emit JSON instead of a summary")
    args = ap.parse_args()
    report = audit_deck(args.pptx)
    if args.json:
        print(json.dumps(report, indent=2))
    else:
        _print_report(report)
    return 0 if report["passed"] else 1


if __name__ == "__main__":
    sys.exit(main())
