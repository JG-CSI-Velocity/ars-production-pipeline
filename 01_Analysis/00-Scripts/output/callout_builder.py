"""Callout-box composition (Phase T1.4 / issue #149).

Replaces the ad-hoc ``_add_callout_box`` body in ``deck_builder.py`` with a
structured ``CalloutBox`` dataclass + ``CalloutBoxBuilder`` so every
content slide carries the same shape: metric + value + denominator +
comparison + (optional) section accent and icon.

SLIDE_DESIGN.md §7 is the authority.

The dataclass is consumed by ``deck_builder._add_slide`` via
``CalloutBoxBuilder.from_kpis(slide_type, section_key, kpis, ...)``.
Existing callers that only pass a raw ``kpis`` dict get a reasonable
default (the first non-subtitle entry becomes ``value`` and ``metric``)
so the refactor is backwards-compatible.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ars_analysis.shared.charts_palette import COLORS, section_color


# ---------------------------------------------------------------------------
# Section -> icon. Plain ASCII so PPTX renders identically on any font set.
# Emoji intentionally avoided -- python-pptx font fallback for emoji is
# inconsistent across Windows / Mac and produces the dreaded tofu square
# in client-facing decks.
# ---------------------------------------------------------------------------
_SECTION_ICONS = {
    "overview":    "▣",
    "dctr":        "◉",
    "rege":        "✦",
    "attrition":   "↘",
    "value":       "$",
    "mailer":      "✉",
    "insights":    "★",
    "transaction": "⇄",
    "ics":         "◆",
}


# Slide types that should NOT receive an auto-callout. screenshot_kpi
# already renders its KPIs natively; title + section are decorative.
_SKIP_TYPES = {"title", "section", "screenshot_kpi", "blank"}


@dataclass
class CalloutBox:
    """Structured payload for a single callout (SLIDE_DESIGN.md §7).

    Required fields:
        metric        — label shown beneath the hero number (e.g. "DCTR")
        value         — the hero number itself, pre-formatted ("34%", "$1.2M")
    Optional context (any may be empty):
        denominator   — population the value is measured against
                        ("of 12,400 eligible accounts")
        comparison    — peer / benchmark / prior-period context
                        ("vs 41% peer median")
        insight       — one-line takeaway for the speaker notes / Excel review
        section_key   — drives accent color + icon (one of SECTION_COLORS keys)
        icon          — override the auto-selected section icon
        position      — "bottom_right" (default) | "bottom_left"
    """

    metric: str
    value: str
    denominator: str = ""
    comparison: str = ""
    insight: str = ""
    section_key: str | None = None
    icon: str | None = None
    position: str = "bottom_right"
    # Filled in by the builder; not part of caller's contract.
    accent_hex: str = field(default="", init=False)

    def __post_init__(self) -> None:
        self.accent_hex = section_color(self.section_key) if self.section_key else COLORS["secondary"]
        if self.icon is None and self.section_key:
            self.icon = _SECTION_ICONS.get(self.section_key, "")
        if self.icon is None:
            self.icon = ""


class CalloutBoxBuilder:
    """Render a :class:`CalloutBox` onto a python-pptx slide.

    Owns the geometry, color, and font rules. The caller (deck_builder)
    just hands it a slide and a CalloutBox.
    """

    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)
    BOX_W = Inches(4.0)
    BOX_H = Inches(1.6)
    FOOTER_CLEARANCE = Inches(0.75)
    LEFT_MARGIN = Inches(0.86)
    RIGHT_MARGIN = Inches(0.5)

    @classmethod
    def from_kpis(
        cls,
        kpis: dict[str, Any] | None,
        *,
        slide_type: str,
        section_key: str | None = None,
        denominator: str = "",
        comparison: str = "",
        insight: str = "",
    ) -> CalloutBox | None:
        """Backwards-compatible factory from the legacy ``kpis`` dict.

        Picks the first non-subtitle entry as the hero ``(metric, value)``;
        any second entry becomes the ``comparison`` line if the caller
        didn't pass one. Returns None if the slide_type should skip a
        callout or the kpis dict yields no hero value.
        """
        if slide_type in _SKIP_TYPES or not kpis:
            return None
        hero_label = hero_value = None
        secondary = None
        for label, value in kpis.items():
            if str(label).lower() in ("subtitle", "title"):
                continue
            if hero_value is None:
                hero_label, hero_value = str(label), str(value)
            elif secondary is None:
                secondary = f"{label}: {value}"
        if hero_value is None:
            return None
        return CalloutBox(
            metric=hero_label or "",
            value=hero_value,
            denominator=denominator,
            comparison=comparison or (secondary or ""),
            insight=insight,
            section_key=section_key,
        )

    @classmethod
    def render(cls, slide, callout: CalloutBox) -> bool:
        """Paint the callout. Returns True on success, False if skipped.

        Skips if no collision-free placement exists (mirrors the picture-
        bounds avoidance from b0ff9d6 / issue 3.7).
        """
        if not callout or not callout.value:
            return False

        # Collision-aware placement -- check existing picture shapes.
        pic_boxes = []
        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            try:
                pic_boxes.append((shape.left, shape.top,
                                  shape.left + shape.width,
                                  shape.top + shape.height))
            except Exception:
                continue

        right_x = cls.SLIDE_W - cls.BOX_W - cls.RIGHT_MARGIN
        left_x = cls.LEFT_MARGIN
        default_top = cls.SLIDE_H - cls.BOX_H - cls.FOOTER_CLEARANCE

        def _overlaps(left, top, right, bottom):
            for pl, pt, pr, pb in pic_boxes:
                if right <= pl or left >= pr:
                    continue
                if bottom <= pt or top >= pb:
                    continue
                return True
            return False

        order = [(left_x, default_top), (right_x, default_top)] if callout.position == "bottom_left" \
                else [(right_x, default_top), (left_x, default_top)]
        chosen = None
        for cand_left, cand_top in order:
            if not _overlaps(cand_left, cand_top, cand_left + cls.BOX_W, cand_top + cls.BOX_H):
                chosen = (cand_left, cand_top, cls.BOX_W, cls.BOX_H)
                break
        if chosen is None and pic_boxes:
            pic_bottom = max(pb for _, _, _, pb in pic_boxes)
            narrow_w = Inches(3.4)
            narrow_h = Inches(1.0)
            narrow_top = pic_bottom + Inches(0.1)
            footer_top = cls.SLIDE_H - Inches(0.5)
            if narrow_top + narrow_h <= footer_top:
                chosen = (cls.SLIDE_W - narrow_w - cls.RIGHT_MARGIN,
                          narrow_top, narrow_w, narrow_h)
        if chosen is None:
            return False

        box_left, box_top, box_w, box_h = chosen
        accent = _hex_to_rgb(callout.accent_hex)

        # Background -- subtle tinted fill so the callout reads as a
        # related-but-distinct region, not a duplicate chart.
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _tint(accent, 0.92)
        shape.line.color.rgb = accent
        shape.line.width = Pt(1.5)
        try:
            shape.adjustments[0] = 0.12
        except Exception:
            pass

        # Top row: icon + metric label, accent-colored.
        if callout.icon:
            _put_text(
                slide,
                box_left + Inches(0.2), box_top + Inches(0.08),
                Inches(0.35), Inches(0.35),
                callout.icon, size=16, bold=True, rgb=accent,
            )
            label_left = box_left + Inches(0.55)
            label_w = box_w - Inches(0.75)
        else:
            label_left = box_left + Inches(0.2)
            label_w = box_w - Inches(0.4)
        _put_text(
            slide, label_left, box_top + Inches(0.10), label_w, Inches(0.3),
            callout.metric.upper(), size=10, bold=True, rgb=accent,
        )

        # Hero number.
        _put_text(
            slide,
            box_left + Inches(0.2), box_top + Inches(0.38),
            box_w - Inches(0.4), Inches(0.55),
            callout.value, size=32, bold=True, rgb=accent,
        )

        # Denominator.
        if callout.denominator:
            _put_text(
                slide,
                box_left + Inches(0.2), box_top + Inches(0.96),
                box_w - Inches(0.4), Inches(0.25),
                callout.denominator, size=11, bold=False, rgb=RGBColor(0x33, 0x33, 0x33),
            )

        # Comparison.
        if callout.comparison:
            _put_text(
                slide,
                box_left + Inches(0.2), box_top + Inches(1.22),
                box_w - Inches(0.4), Inches(0.25),
                callout.comparison, size=10, bold=False,
                rgb=RGBColor(0x55, 0x55, 0x55), italic=True,
            )

        return True


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _hex_to_rgb(hexstr: str) -> RGBColor:
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _tint(rgb: RGBColor, amount: float) -> RGBColor:
    """Mix ``rgb`` with white by ``amount`` (0=pure rgb, 1=white)."""
    r, g, b = rgb[0], rgb[1], rgb[2]
    r2 = int(r + (255 - r) * amount)
    g2 = int(g + (255 - g) * amount)
    b2 = int(b + (255 - b) * amount)
    return RGBColor(r2, g2, b2)


def _put_text(slide, left, top, width, height, text, *,
              size: int, bold: bool, rgb: RGBColor,
              italic: bool = False, align=PP_ALIGN.LEFT) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = rgb
    p.alignment = align
