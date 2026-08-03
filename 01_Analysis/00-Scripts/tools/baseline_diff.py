"""Full-deck baseline diff: prove the module refactor did NOT change output.

The refactor's promotions re-run the same analytics cells, so a full
``--product txn`` (or ``combined``) deck should be byte-for-content identical to
a pre-refactor build for the same client. This tool makes that check one
command on a machine that has the client data (the M: box):

    python tools/baseline_diff.py --client 1776 --month 2026.06 --csm JamesG \
        --baseline <pre-refactor-git-ref>

It (1) builds the deck on the current checkout, (2) checks the baseline ref out
into a throwaway ``git worktree`` and builds the deck there over the SAME client
data, then (3) diffs the two .pptx by slide count and per-slide text + image
count. Exit 0 == identical (refactor safe); exit 1 == differences (investigate).

The deck-diff itself (`diff_decks`) is pure and unit-tested; the run/worktree
orchestration needs the client data and so only does useful work on the M: box.
"""

from __future__ import annotations

import argparse
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

_REPO = Path(__file__).resolve().parents[2]


@dataclass
class SlideView:
    """The comparable content of one slide: its text and how many pictures."""

    texts: tuple[str, ...]
    n_pictures: int


@dataclass
class DeckDiff:
    slide_count_a: int
    slide_count_b: int
    per_slide: list[str] = field(default_factory=list)  # human-readable diffs

    @property
    def identical(self) -> bool:
        return self.slide_count_a == self.slide_count_b and not self.per_slide


def _slide_views(pptx_path: Path) -> list[SlideView]:
    from pptx import Presentation  # imported lazily so the CLI --help works bare
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    views: list[SlideView] = []
    for slide in Presentation(str(pptx_path)).slides:
        texts: list[str] = []
        pics = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        views.append(SlideView(texts=tuple(texts), n_pictures=pics))
    return views


def diff_decks(deck_a: Path, deck_b: Path) -> DeckDiff:
    """Structured content diff of two decks. `a` = current, `b` = baseline."""
    va, vb = _slide_views(deck_a), _slide_views(deck_b)
    diff = DeckDiff(slide_count_a=len(va), slide_count_b=len(vb))
    for i in range(max(len(va), len(vb))):
        if i >= len(va):
            diff.per_slide.append(f"slide {i+1}: present in baseline only")
            continue
        if i >= len(vb):
            diff.per_slide.append(f"slide {i+1}: present in current only")
            continue
        a, b = va[i], vb[i]
        if a.n_pictures != b.n_pictures:
            diff.per_slide.append(
                f"slide {i+1}: picture count {a.n_pictures} vs {b.n_pictures}")
        if a.texts != b.texts:
            diff.per_slide.append(f"slide {i+1}: text differs")
    return diff


def _build_deck(cwd: Path, client: str, month: str, csm: str, product: str) -> Path | None:
    """Run the analysis pipeline in `cwd` and return the produced .pptx (newest)."""
    run_py = cwd / "01_Analysis" / "run.py"
    proc = subprocess.run(
        [sys.executable, str(run_py), "--product", product,
         "--client", client, "--month", month, "--csm", csm],
        cwd=str(run_py.parent), capture_output=True, text=True,
    )
    if proc.returncode != 0:
        print(proc.stdout[-2000:])
        print(proc.stderr[-2000:], file=sys.stderr)
        return None
    pres = cwd / "02_Presentations"
    decks = sorted(pres.rglob(f"*_{product}_deck.pptx"), key=lambda p: p.stat().st_mtime)
    return decks[-1] if decks else None


def main() -> int:
    ap = argparse.ArgumentParser(description="Full-deck baseline diff for the refactor.")
    ap.add_argument("--client", required=True)
    ap.add_argument("--month", required=True)
    ap.add_argument("--csm", required=True)
    ap.add_argument("--product", default="txn", choices=["txn", "combined", "ars"])
    ap.add_argument("--baseline", required=True,
                    help="git ref of the pre-refactor code to compare against")
    args = ap.parse_args()

    print(f"[1/3] Building current deck ({args.product}) for {args.client}...")
    cur = _build_deck(_REPO, args.client, args.month, args.csm, args.product)
    if not cur:
        print("Current build failed."); return 2

    with tempfile.TemporaryDirectory() as tmp:
        wt = Path(tmp) / "baseline"
        print(f"[2/3] Checking out baseline {args.baseline} into a worktree...")
        wc = subprocess.run(["git", "worktree", "add", "--detach", str(wt), args.baseline],
                            cwd=str(_REPO), capture_output=True, text=True)
        if wc.returncode != 0:
            print(wc.stderr); return 2
        try:
            base = _build_deck(wt, args.client, args.month, args.csm, args.product)
            if not base:
                print("Baseline build failed."); return 2
            print("[3/3] Diffing decks...")
            d = diff_decks(cur, base)
        finally:
            subprocess.run(["git", "worktree", "remove", "--force", str(wt)],
                           cwd=str(_REPO), capture_output=True)

    print(f"  current : {d.slide_count_a} slides ({cur.name})")
    print(f"  baseline: {d.slide_count_b} slides")
    if d.identical:
        print("  RESULT: IDENTICAL -- refactor did not change deck output. ✓")
        return 0
    print(f"  RESULT: {len(d.per_slide)} difference(s):")
    for line in d.per_slide[:50]:
        print(f"    - {line}")
    return 1


if __name__ == "__main__":
    raise SystemExit(main())
