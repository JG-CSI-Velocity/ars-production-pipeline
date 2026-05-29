"""
Phase 18.1 / 18.2 / 18.3 — deck_builder.py patches
====================================================

File: 01_Analysis/00-Scripts/output/deck_builder.py

These are drop-in methods and modifications for the existing DeckBuilder class.
Apply them in order. Each section has:
  - WHERE: exact location in the existing file
  - WHAT: the new/replaced code
  - WHY: which SLIDE_DESIGN.md section it implements

Tested against: SlideContent dataclass, all LAYOUT_* constants, existing
_add_fitted_picture, _set_title, _get_image_positioning helpers.
"""

# =============================================================================
# IMPORTS TO ADD (at top of deck_builder.py, near existing pptx imports)
# =============================================================================
# These may already be imported — add only what's missing:
#
#   from pptx.util import Emu
#   from pptx.oxml.ns import nsmap
#
# Emu is needed for precise shape positioning in footer bands.
# nsmap is already imported (used in build() for qn()).


# =============================================================================
# 18.1 — CALLOUT BOX RENDERING
# =============================================================================
# SLIDE_DESIGN.md §7: Every content slide has one callout.
#   - Hero number: accent color, 32pt bold
#   - Sub-label: 14pt semibold + 12pt regular (2 lines max)
#   - Anchored to a specific action the client can take
#
# WHERE: Add as a new method on DeckBuilder class, after _add_fitted_picture()
# CALLED FROM: _build_screenshot_slide, _build_multi_screenshot_slide
# (screenshot_kpi already has its own KPI rendering — skip for that type)

def _add_callout_box(self, slide, kpis, position="bottom_right"):
    """Add a callout box with hero number + sub-label per SLIDE_DESIGN.md §7.

    Args:
        slide: The pptx slide object.
        kpis: dict of label->value pairs from AnalysisResult.kpis.
              First non-'subtitle' entry becomes the hero number.
        position: 'bottom_right' (default) or 'bottom_left'.
    """
    if not kpis:
        return

    # Extract hero KPI (first non-subtitle entry)
    hero_label = None
    hero_value = None
    sub_label = None
    for label, value in kpis.items():
        if label.lower() in ("subtitle", "title"):
            continue
        if hero_value is None:
            hero_label = label
            hero_value = str(value)
        elif sub_label is None:
            sub_label = f"{label}: {value}"

    if hero_value is None:
        return

    # Callout box dimensions
    box_width = Inches(3.8)
    box_height = Inches(1.4)

    if position == "bottom_left":
        box_left = Inches(0.86)
    else:
        # bottom_right: anchored to right margin
        box_left = Inches(13.33) - box_width - Inches(0.5)

    box_top = Inches(7.5) - box_height - Inches(0.75)  # above footer band

    # Draw rounded rectangle background
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        box_left, box_top, box_width, box_height,
    )
    # Light teal fill, subtle border
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xED, 0xFA, 0xF9)  # very light teal
    shape.line.color.rgb = RGBColor(0x0D, 0x94, 0x88)        # teal border
    shape.line.width = Pt(1.5)

    # Round the corners (adjustment value 0-100000, higher = rounder)
    try:
        shape.adjustments[0] = 0.15
    except Exception:
        pass

    # Hero number — large, bold, teal
    tb_value = slide.shapes.add_textbox(
        box_left + Inches(0.2),
        box_top + Inches(0.1),
        box_width - Inches(0.4),
        Inches(0.6),
    )
    tf = tb_value.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = hero_value
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0D, 0x94, 0x88)  # teal accent
    p.alignment = PP_ALIGN.LEFT

    # Sub-label line 1: hero label
    tb_label = slide.shapes.add_textbox(
        box_left + Inches(0.2),
        box_top + Inches(0.7),
        box_width - Inches(0.4),
        Inches(0.3),
    )
    tf = tb_label.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = hero_label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x3D, 0x59)  # navy
    p.alignment = PP_ALIGN.LEFT

    # Sub-label line 2: secondary KPI (if present)
    if sub_label:
        tb_sub = slide.shapes.add_textbox(
            box_left + Inches(0.2),
            box_top + Inches(1.0),
            box_width - Inches(0.4),
            Inches(0.3),
        )
        tf = tb_sub.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sub_label
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.alignment = PP_ALIGN.LEFT


# =============================================================================
# 18.2 — FOOTER BAND (EVERY SLIDE)
# =============================================================================
# SLIDE_DESIGN.md §8:
#   Line 1: Source/methodology — 9pt italic, #777777
#   Line 2: {Client Name} | {Month YYYY} | Slide {n} | STRICTLY CONFIDENTIAL — 8pt, #999999
#
# WHERE: Add as new method on DeckBuilder class.
# CALLED FROM: _add_slide(), AFTER the builder dispatch, BEFORE notes.
# Skipped for: title slides, section dividers (they have their own styling).

def _add_footer_band(self, slide, content, slide_number, client_name="", month=""):
    """Add two-line footer band per SLIDE_DESIGN.md §8.

    Args:
        slide: The pptx slide object.
        content: SlideContent for this slide.
        slide_number: 1-indexed slide number in the deck.
        client_name: Client display name.
        month: Month string (YYYY.MM format, converted to display).
    """
    # Skip footer on title and section slides
    if content.slide_type in ("title", "section"):
        return

    slide_width = Inches(13.33)
    footer_left = Inches(0.5)
    footer_width = slide_width - Inches(1.0)

    # Line 1: Source/methodology (from notes or default)
    source_text = ""
    if content.notes_text:
        # Extract first line of notes as source hint
        first_line = content.notes_text.split("\n")[0]
        if first_line.startswith("KEY FINDING:"):
            source_text = ""  # Don't repeat the headline
        else:
            source_text = first_line[:120]

    if not source_text:
        source_text = "Source: Client ODD file | Analysis by Velocity Solutions"

    tb_source = slide.shapes.add_textbox(
        footer_left, Inches(7.0), footer_width, Inches(0.2),
    )
    tf = tb_source.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = source_text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
    p.alignment = PP_ALIGN.LEFT

    # Line 2: Client | Month | Slide # | STRICTLY CONFIDENTIAL
    # Convert YYYY.MM to "Month YYYY"
    month_display = month
    if month and "." in month:
        try:
            import calendar as _cal
            parts = month.split(".")
            month_display = f"{_cal.month_name[int(parts[1])]} {parts[0]}"
        except (ValueError, IndexError):
            month_display = month

    footer_parts = []
    if client_name:
        footer_parts.append(client_name)
    if month_display:
        footer_parts.append(month_display)
    footer_parts.append(f"Slide {slide_number}")
    footer_parts.append("STRICTLY CONFIDENTIAL")

    footer_text = "  |  ".join(footer_parts)

    tb_footer = slide.shapes.add_textbox(
        footer_left, Inches(7.22), footer_width, Inches(0.2),
    )
    tf = tb_footer.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.LEFT


# =============================================================================
# 18.3 — SECTION DIVIDERS
# =============================================================================
# SLIDE_DESIGN.md §9:
#   - Full-bleed navy background
#   - Section number in teal: "02"
#   - Section title: white, 36pt bold
#   - Lead-in sentence: white, 18pt regular
#   - No charts, no logos, no ornamentation
#
# WHERE: REPLACE existing _build_section_slide method on DeckBuilder class.
#
# The section number + lead-in are encoded in the SlideContent.title field.
# Convention: "Section Number\nSection Title\nLead-in sentence"
# If only "Title\nSubtitle" is provided, we skip the section number.

def _build_section_slide(self, slide, content):
    """Build section divider slide per SLIDE_DESIGN.md §9.

    Expected title format (any of):
      "02\nDebit Card Take Rate\nCurrent penetration and opportunity"
      "Debit Card Take Rate\nSubtitle text"
      "Debit Card Take Rate"
    """
    # Remove all placeholders — we draw everything programmatically
    for ph in slide.placeholders:
        try:
            ph.element.getparent().remove(ph.element)
        except Exception:
            pass

    # Set slide background to navy (#1B365D)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x36, 0x5D)

    # Parse title into components
    lines = content.title.split("\n") if content.title else [""]
    section_number = None
    section_title = lines[0]
    lead_in = None

    if len(lines) >= 3:
        # Full format: number, title, lead-in
        # Check if first line looks like a section number (1-2 digits or "0X")
        first = lines[0].strip()
        if len(first) <= 3 and (first.isdigit() or first.startswith("0")):
            section_number = first
            section_title = lines[1]
            lead_in = lines[2]
        else:
            section_title = lines[0]
            lead_in = "\n".join(lines[1:])
    elif len(lines) == 2:
        section_title = lines[0]
        lead_in = lines[1]

    # Y-positioning: centered vertically with breathing room
    y_cursor = Inches(2.5)

    # Section number (teal, 48pt bold)
    if section_number:
        tb = slide.shapes.add_textbox(
            Inches(0.86), Inches(2.2), Inches(2.0), Inches(0.7),
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = section_number
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0D, 0x94, 0x88)  # teal
        p.alignment = PP_ALIGN.LEFT
        y_cursor = Inches(3.0)

    # Section title (white, 36pt bold)
    tb = slide.shapes.add_textbox(
        Inches(0.86), y_cursor, Inches(10.0), Inches(1.0),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT

    # Lead-in sentence (white, 18pt regular)
    if lead_in:
        tb = slide.shapes.add_textbox(
            Inches(0.86), y_cursor + Inches(1.1), Inches(10.0), Inches(0.8),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = lead_in
        p.font.size = Pt(18)
        p.font.bold = False
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.LEFT


# =============================================================================
# MODIFIED _add_slide — Wire callout + footer into the dispatch
# =============================================================================
# WHERE: REPLACE existing _add_slide method on DeckBuilder class.
#
# Changes vs existing:
#   1. Track slide_number (1-indexed)
#   2. After builder dispatch, call _add_callout_box for screenshot types
#   3. After builder dispatch, call _add_footer_band for content slides
#   4. Pass client_name/month through (requires build() to store them on self)

def _add_slide(self, content, slide_number=0, client_name="", month=""):
    """Create single slide and dispatch to appropriate builder method.

    Enhanced: adds callout boxes (§7) and footer bands (§8) after content.
    """
    slide = self.prs.slides.add_slide(self.prs.slide_layouts[content.layout_index])

    builders = {
        "title": self._build_title_slide,
        "section": self._build_section_slide,
        "screenshot": self._build_screenshot_slide,
        "screenshot_kpi": self._build_screenshot_kpi_slide,
        "multi_screenshot": self._build_multi_screenshot_slide,
        "summary": self._build_summary_slide,
        "blank": self._build_blank_slide,
        "mailer_summary": self._build_mailer_summary_slide,
        "kpi_dashboard": self._build_kpi_dashboard_slide,
        "chart_narrative": self._build_chart_narrative_slide,
        "kpi_hero": self._build_kpi_hero_slide,
    }

    builder = builders.get(content.slide_type)
    if builder:
        builder(slide, content)
    else:
        logger.warning("Unknown slide_type: {t}", t=content.slide_type)

    # --- Phase 18.1: Callout box on chart slides ---
    # screenshot_kpi already renders its own KPIs; skip to avoid double-render.
    # Also skip title, section, blank, summary — they don't have chart areas.
    if content.slide_type in ("screenshot", "multi_screenshot") and content.kpis:
        self._add_callout_box(slide, content.kpis)

    # --- Phase 18.2: Footer band on content slides ---
    if slide_number > 0:
        self._add_footer_band(slide, content, slide_number, client_name, month)

    # Write speaker notes if provided
    if content.notes_text:
        try:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                tf.text = content.notes_text
        except Exception:
            pass


# =============================================================================
# MODIFIED build() — Pass client_name, month, slide_number through
# =============================================================================
# WHERE: REPLACE existing build() method on DeckBuilder class.
#
# Changes vs existing:
#   1. Accept client_name and month parameters
#   2. Pass slide_number to _add_slide for footer rendering

def build(self, slides, output_path, client_name="", month=""):
    """Build complete PowerPoint deck from slide definitions.

    Args:
        slides: list of SlideContent objects.
        output_path: Where to save the .pptx file.
        client_name: Client display name (for footer bands).
        month: Month string YYYY.MM (for footer bands).

    Returns:
        output_path string.
    """
    self.prs = Presentation(self.template_path)

    # Remove sample slides that ship with the 2025 template (18 slides)
    while len(self.prs.slides) > 0:
        rId = self.prs.slides._sldIdLst[0].get(qn("r:id"))
        self.prs.part.drop_rel(rId)
        self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    n_layouts = len(self.prs.slide_layouts)

    for i, slide_content in enumerate(slides):
        if slide_content.layout_index >= n_layouts:
            logger.warning(
                "Slide {i} '{title}' has layout_index={idx} but template only has {n} layouts, using 0",
                i=i,
                title=slide_content.title[:40],
                idx=slide_content.layout_index,
                n=n_layouts,
            )
            slide_content.layout_index = 0
        try:
            self._add_slide(
                slide_content,
                slide_number=i + 1,
                client_name=client_name,
                month=month,
            )
        except Exception as exc:
            logger.error(
                "Slide {i} '{title}' (type={t}, layout={l}) failed: {err}",
                i=i,
                title=slide_content.title[:40],
                t=slide_content.slide_type,
                l=slide_content.layout_index,
                err=exc,
            )

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    self.prs.save(output_path)
    return output_path


# =============================================================================
# CALLER UPDATE: build_deck() must pass client_name + month to DeckBuilder.build()
# =============================================================================
# In the existing build_deck(ctx) function (module-level, not on DeckBuilder),
# find the line that calls:
#
#     result = builder.build(all_slides, str(pptx_path))
#
# And change it to:
#
#     result = builder.build(
#         all_slides,
#         str(pptx_path),
#         client_name=ctx.client.client_name,
#         month=ctx.client.month,
#     )
#
# This threads client_name and month from PipelineContext through to the
# footer band renderer.


# =============================================================================
# SECTION DIVIDER TITLE FORMAT UPDATE
# =============================================================================
# The existing code that creates section divider SlideContent objects
# (in build_deck or _build_preamble_slides) currently uses:
#
#     SlideContent(
#         slide_type="section",
#         title=f"SECTION: {label}\n{subtitle}",
#         layout_index=LAYOUT_SECTION_ALT,
#     )
#
# Update these to use the new 3-line format with section numbers:
#
#     _SECTION_NUMBERS = {
#         "overview": "01",
#         "dctr": "02",
#         "rege": "03",
#         "attrition": "04",
#         "mailer": "05",
#         "value": "06",
#         "insights": "07",
#         "competition": "08",
#         "ics": "09",
#     }
#
#     _SECTION_LEAD_INS = {
#         "overview": "Portfolio composition, eligibility, and program scope.",
#         "dctr": "Current penetration, where opportunity sits, and what closing the gap is worth.",
#         "rege": "Opt-in rates, branch variation, and revenue impact of Reg E enrollment.",
#         "attrition": "Who is leaving, what it costs, and where retention efforts should focus.",
#         "mailer": "Campaign response rates, cohort lift, and mailer program ROI.",
#         "value": "Revenue attribution — what a debit card and Reg E opt-in are worth.",
#         "insights": "Synthesis, recommendations, and quantified action plan.",
#         "competition": "Competitor detection, wallet share analysis, and market positioning.",
#         "ics": "Invitation Checking System acquisition channels and conversion.",
#     }
#
# Then when building section dividers:
#
#     section_num = _SECTION_NUMBERS.get(section_key, "")
#     lead_in = _SECTION_LEAD_INS.get(section_key, "")
#     label = _SECTION_LABELS.get(section_key, section_key.title())
#
#     SlideContent(
#         slide_type="section",
#         title=f"{section_num}\n{label}\n{lead_in}",
#         layout_index=LAYOUT_SECTION,  # any section layout works; bg is overridden
#     )
